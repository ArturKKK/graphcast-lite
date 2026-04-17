"""Generate presentation PPTX from slide content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x0B, 0x80, 0x43)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HEADER_BG = RGBColor(0x1A, 0x73, 0xE8)
TABLE_ALT_BG = RGBColor(0xF0, 0xF4, 0xF9)

def add_slide(title_text=None):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide

def add_title(slide, text, top=Inches(0.3), font_size=36, bold=True):
    left = Inches(0.7)
    width = prs.slide_width - Inches(1.4)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = BLACK
    return txBox

def add_text(slide, text, top, left=Inches(0.7), width=None, font_size=20, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    if width is None:
        width = prs.slide_width - Inches(1.4)
    height = Inches(5.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Handle bold markers
        if line.startswith('**') and ':**' in line:
            parts = line.split(':**', 1)
            label = parts[0].replace('**', '')
            rest = parts[1].strip() if len(parts) > 1 else ''
            run1 = p.add_run()
            run1.text = label + ': '
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            if rest:
                run2 = p.add_run()
                run2.text = rest
                run2.font.size = Pt(font_size)
                run2.font.bold = False
                run2.font.color.rgb = color
        else:
            clean = line.replace('**', '')
            is_bold = '**' in line and bold is False
            p.text = clean
            p.font.size = Pt(font_size)
            p.font.bold = bold or is_bold
            p.font.color.rgb = color
        p.alignment = alignment
    return txBox

def add_bullet_list(slide, items, top, left=Inches(0.7), width=None, font_size=20):
    if width is None:
        width = prs.slide_width - Inches(1.4)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Parse bold prefix
        if item.startswith('**') and ':**' in item:
            parts = item.split(':**', 1)
            label = parts[0].replace('**', '')
            rest = parts[1].strip() if len(parts) > 1 else ''
            run1 = p.add_run()
            run1.text = label + ': '
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = DARK
            if rest:
                run2 = p.add_run()
                run2.text = rest
                run2.font.size = Pt(font_size)
                run2.font.bold = False
                run2.font.color.rgb = DARK
        else:
            clean = item.replace('**', '')
            p.text = clean
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK
        p.space_before = Pt(6)
        p.level = 0
        # Indent
        p.space_before = Pt(8)

    return txBox

def add_table(slide, headers, rows, top, left=Inches(0.7), col_widths=None, font_size=16):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        total_w = prs.slide_width - Inches(1.4)
        col_w = int(total_w / n_cols)
        col_widths = [col_w] * n_cols

    table_width = sum(col_widths)
    table_height = Inches(0.4 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, table_height)
    table = shape.table

    for j, w in enumerate(col_widths):
        table.columns[j].width = w

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG

    # Data
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            clean = str(val).replace('**', '')
            cell.text = clean
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = DARK
                paragraph.alignment = PP_ALIGN.CENTER
                if '**' in str(val):
                    paragraph.font.bold = True
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG

    return shape

def add_note(slide, text, top, font_size=16):
    left = Inches(0.9)
    width = prs.slide_width - Inches(1.8)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.italic = True
    p.font.color.rgb = GRAY
    return txBox

def add_placeholder_box(slide, text, top, left=Inches(0.7)):
    width = prs.slide_width - Inches(1.4)
    height = Inches(2.5)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RED
    p.alignment = PP_ALIGN.CENTER
    # Add border via shape line
    shape.line.color.rgb = RED
    shape.line.width = Pt(2)
    shape.line.dash_style = 2  # dash
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
    return shape

def add_formula(slide, text, top, left=Inches(0.7), font_size=18):
    width = prs.slide_width - Inches(1.4)
    height = Inches(0.6)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.CENTER
    return txBox


# ============================================================
# SLIDE 1: Title
# ============================================================
s = add_slide()
add_text(s, "Повышение точности регионального\nпрогноза погоды в графовой нейросетевой\nмодели с помощью усвоения данных.",
         top=Inches(1.5), font_size=40, bold=True, color=BLACK, alignment=PP_ALIGN.LEFT)
add_text(s, "Подготовил: Табаков Артур Станиславович\nНаучный руководитель: Пененко Алексей Владимирович, д.ф.-м.н., ВНС ИВМиМГ СО РАН",
         top=Inches(4.5), font_size=22, color=DARK)

# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
s = add_slide()
add_title(s, "Постановка проблемы")
add_bullet_list(s, [
    "**Накопление ошибки:** В авторегрессионном режиме нейросеть быстро теряет точность, если не корректировать её траекторию.",
    "**Разреженность данных:** Реальная сеть метеостанций покрывает лишь малую часть территории. Нам нужно уметь восстанавливать прогноз в «слепых зонах».",
    "**Граничные условия:** При вырезании региона из глобальной сетки возникают искусственные границы, где отсутствие информации о соседях приводит к численным артефактам.",
], top=Inches(1.4), font_size=22)

# ============================================================
# SLIDE 3: Goals
# ============================================================
s = add_slide()
add_title(s, "Цель и задачи")
add_text(s, "Цель: Повысить точность краткосрочного прогноза погоды над Красноярском\nв рамках графовой нейросетевой модели.",
         top=Inches(1.3), font_size=22, bold=True, color=BLACK)
add_bullet_list(s, [
    "1. Реализовать глобальную GNN-модель прогноза погоды на основе архитектуры GraphCast.",
    "2. Разработать мультирезолюционный подход: встраивание регионального графа высокого разрешения в глобальный.",
    "3. Предложить стратегию дообучения с заморозкой процессора для сохранения глобальных знаний.",
    "4. Реализовать методы усвоения данных (Nudging, Optimal Interpolation) для коррекции по разреженным наблюдениям.",
    "5. Исследовать влияние различных групп переменных на качество усвоения.",
], top=Inches(2.5), font_size=20)

# ============================================================
# SLIDE 4: Architecture
# ============================================================
s = add_slide()
add_title(s, "Базовая модель: Архитектура GraphCast")
add_bullet_list(s, [
    "**Encoder:** Проекция данных с регулярной сетки (grid) на икосаэдральную (mesh).",
    "**Processor:** Обновление состояний на mesh с помощью Message Passing (12 шагов).",
    "**Decoder:** Интерполяция прогноза обратно на регулярную сетку.",
], top=Inches(1.3), font_size=20)
add_text(s, "Processor построен на базе Interaction Network (Battaglia et al., 2016):",
         top=Inches(3.2), font_size=20, bold=True)
add_bullet_list(s, [
    "На каждом шаге MP рёбра и узлы обновляются отдельными MLP.",
    "Edge MLP: фичи двух узлов + фичи ребра → обновлённое сообщение.",
    "Node MLP: фичи узла + агрегат входящих сообщений → обновлённое состояние.",
    "Рёбра — полноценные объекты с собственными обучаемыми представлениями.",
], top=Inches(3.9), font_size=18)
add_text(s, "Параметры: dim=256, edge_dim=4, mesh levels [4, 6].",
         top=Inches(5.8), font_size=18, color=GRAY)
add_note(s, "Выбор GNN обусловлен сферической геометрией Земли для устранения искажений на полюсах.",
         top=Inches(6.4))

# ============================================================
# SLIDE 5: Metrics
# ============================================================
s = add_slide()
add_title(s, "Метрики оценки")
add_text(s, "Реанализ ERA5 — «лучшая оценка» атмосферы: ECMWF объединяет спутниковые, станционные и зондовые данные\nчерез физическую модель с усвоением → полный 3D-снимок атмосферы. Используется как ground truth.",
         top=Inches(1.2), font_size=18, color=GRAY)
add_bullet_list(s, [
    "**RMSE:** среднеквадратичная ошибка (в физ. единицах: °C, м/с, гПа).",
    "**RMSE (норм.):** RMSE на нормализованных данных (z-score). Безразмерная, для сравнения по всем переменным.",
    "**Persistence:** наивный бейзлайн «завтра = сегодня». Любая полезная модель должна быть лучше.",
    "**Skill Score:** Skill = 1 − RMSE_model / RMSE_persistence. Насколько модель лучше persistence (в %).",
    "**ACC:** корреляция предсказанных и реальных аномалий. ACC=1 — идеал, ACC≤0 — не лучше случайного.",
], top=Inches(2.2), font_size=20)

# ============================================================
# SLIDE 5.5: Variables
# ============================================================
s = add_slide()
add_title(s, "Набор переменных (19 полей)")
add_table(s,
    ["Переменная", "Описание", "Уровень"],
    [
        ["t2m", "Температура воздуха на 2м", "поверхность"],
        ["10u, 10v", "Зональная / меридиональная компонента ветра", "10м"],
        ["msl", "Давление, приведённое к уровню моря", "поверхность"],
        ["sp", "Давление на поверхности", "поверхность"],
        ["tp", "Суммарные осадки", "поверхность"],
        ["tcwv", "Общее влагосодержание столба атмосферы", "интегральная"],
        ["z_surf", "Геопотенциал поверхности (рельеф)", "статическое"],
        ["lsm", "Маска суша/море", "статическое"],
        ["t, u, v, z, q @850", "Температура, ветер, геопотенциал, влажность", "850 гПа (~1.5 км)"],
        ["t, u, v, z, q @500", "Температура, ветер, геопотенциал, влажность", "500 гПа (~5.5 км)"],
    ],
    top=Inches(1.3), left=Inches(0.7),
    col_widths=[Inches(2.5), Inches(6), Inches(3)],
    font_size=16)
add_note(s, "Данные: реанализ ERA5 (ECMWF), период 2010–2021, шаг 6 часов.", top=Inches(6.2))

# ============================================================
# SLIDE 6: Global model results
# ============================================================
s = add_slide()
add_title(s, "Глобальный baseline: прогноз на 3 суток")
add_text(s, "Сетка 512×256 (≈ 0.7°), 19 переменных, 200 тестовых сэмплов, 12 AR-шагов до +72ч.",
         top=Inches(1.1), font_size=18)
add_text(s, "Глобально по всей сетке:", top=Inches(1.55), font_size=18, bold=True)
add_table(s,
    ["Горизонт", "RMSE (норм.)", "Skill vs persistence", "ACC"],
    [
        ["+24ч", "0.245", "**62.3%**", "0.969"],
        ["+48ч", "0.378", "**50.2%**", "0.928"],
        ["+72ч", "0.518", "**34.3%**", "0.869"],
        ["Среднее +6…+72ч", "0.349", "**49.7%**", "0.939"],
    ],
    top=Inches(1.95), left=Inches(1.35),
    col_widths=[Inches(2), Inches(2.5), Inches(2.5), Inches(2)],
    font_size=17)
add_text(s, "В регионе Красноярска на грубой глобальной сетке:", top=Inches(4.35), font_size=18, bold=True)
add_table(s,
    ["Горизонт", "t2m RMSE", "Skill", "ACC"],
    [
        ["+24ч", "0.72°C", "71.5%", "0.860"],
        ["+48ч", "1.10°C", "57.0%", "0.744"],
        ["+72ч", "1.52°C", "41.5%", "0.565"],
    ],
    top=Inches(4.75), left=Inches(2.0),
    col_widths=[Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)],
    font_size=17)
add_note(s, "Важно: на глобальной сетке регион покрыт только 3 узлами, поэтому эти локальные метрики нельзя напрямую сравнивать с multires (45 узлов).", top=Inches(6.35))

# ============================================================
# SLIDE 7: Multires approach
# ============================================================
s = add_slide()
add_title(s, "Мультирезолюционный подход")
add_text(s, "Проблема: Если вырезать регион из глобальной сетки — возникают искусственные границы → артефакты.",
         top=Inches(1.3), font_size=20, bold=True, color=RED)
add_text(s, "Решение: Встраиваем региональный граф высокого разрешения прямо в глобальный,\nформируя единую мультирезолюционную сетку.",
         top=Inches(2.0), font_size=20, bold=True, color=BLACK)
add_bullet_list(s, [
    "Глобальная подсетка: 512×256 = 131 072 узла (шаг ~0.7°)",
    "Региональная подсетка (Красноярск): 61×41 = 2 501 узел (шаг 0.25°)",
    "Итого: 133 279 узлов, единый граф",
], top=Inches(3.0), font_size=20)
add_text(s, "Преимущества:", top=Inches(4.5), font_size=20, bold=True)
add_bullet_list(s, [
    "Нет искусственных границ — глобальные узлы обеспечивают естественный контекст.",
    "Региональные узлы дают локальную детализацию в 3× выше.",
    "Encoder/Decoder работают с объединённой сеткой, Processor — на единой mesh.",
], top=Inches(5.0), font_size=18)
add_note(s, "Целевой регион (~55°–57°N, ~90°–95°E): 45 внутренних узлов вокруг Красноярска — зона оценки метрик.",
         top=Inches(6.5))

# ============================================================
# SLIDE 8: Training strategy
# ============================================================
s = add_slide()
add_title(s, "Стратегия дообучения: заморозка процессора")
add_text(s, "Проблема: При fine-tune на мультирезолюционных данных процессор может «забыть»\nглобальную физику (catastrophic forgetting).",
         top=Inches(1.3), font_size=20, color=RED)
add_text(s, "Решение: Стратегия freeze → fine-tune:", top=Inches(2.2), font_size=20, bold=True)
add_bullet_list(s, [
    "Этап 1 (6 эпох): Processor ЗАМОРОЖЕН. Обучаются только Encoder и Decoder → модель адаптирует проекцию данных, не теряя глобальных знаний.",
    "Этап 2 (оставшиеся эпохи): Processor разморожен с пониженным LR (×0.1) → тонкая подстройка.",
], top=Inches(2.8), font_size=20)
add_table(s,
    ["Параметр", "Значение"],
    [
        ["freeze_epochs", "6"],
        ["lr (Encoder/Decoder)", "1e-4"],
        ["lr (Processor, этап 2)", "1e-5"],
        ["Всего эпох", "32"],
        ["AR шагов при обучении", "4"],
    ],
    top=Inches(4.5), left=Inches(3.5),
    col_widths=[Inches(3), Inches(2.5)],
    font_size=18)

# ============================================================
# SLIDE 9: Baseline evaluation
# ============================================================
s = add_slide()
add_title(s, "Оценка модели (без усвоения данных)")
add_text(s, "Сравнение двух стратегий. 200 тестовых сэмплов (каждый — пара наблюдений → прогноз 4 шага; подвыборка из 1750):",
         top=Inches(1.3), font_size=20)
add_table(s,
    ["Метрика", "freeze6", "nofreeze"],
    [
        ["Skill (глобально)", "**66.9%**", "65.2%"],
        ["Skill (регион Красноярска)", "**75.8%**", "74.5%"],
        ["ACC (глобально)", "**0.983**", "0.981"],
    ],
    top=Inches(2.2), left=Inches(2.5),
    col_widths=[Inches(3.5), Inches(2), Inches(2)],
    font_size=18)
add_text(s, "Региональная RMSE по t2m (°C):", top=Inches(4.0), font_size=20, bold=True)
add_table(s,
    ["Горизонт", "freeze6", "nofreeze", "Δ"],
    [
        ["+6ч", "0.96", "0.98", "−0.02"],
        ["+12ч", "1.22", "1.33", "−0.11"],
        ["+18ч", "1.29", "1.60", "−0.31"],
        ["+24ч", "1.40", "1.82", "−0.42"],
    ],
    top=Inches(4.6), left=Inches(2.5),
    col_widths=[Inches(2), Inches(1.8), Inches(1.8), Inches(1.5)],
    font_size=18)
add_note(s, "Вывод: Заморозка процессора критична — без неё ошибка на +24ч растёт на 30%.", top=Inches(6.5))

# ============================================================
# SLIDE 10: Nudging
# ============================================================
s = add_slide()
add_title(s, "Усвоение данных: Nudging")
add_text(s, "Задача: Скорректировать прогноз, имея наблюдения лишь в 10% узлов сетки\n(имитация разреженной сети станций).",
         top=Inches(1.3), font_size=20)
add_text(s, "Принцип: Эвристическая итеративная коррекция — «подтягиваем» прогноз к наблюдению:",
         top=Inches(2.2), font_size=20, bold=True)
add_formula(s, "xₐ = x_b + α · M · (y_obs − H(x_b))", top=Inches(3.0), font_size=24)
add_bullet_list(s, [
    "x_b — фоновый прогноз",
    "y_obs — вектор наблюдений",
    "M — маска доступности данных (1 в точках станций, 0 в остальных)",
    "α — коэффициент релаксации (оптимально α = 0.5)",
], top=Inches(3.8), font_size=18)
add_note(s, "Nudging выполняется перед подачей данных в модель на каждом авторегрессионном шаге.", top=Inches(5.8))

# ============================================================
# SLIDE 11: OI
# ============================================================
s = add_slide()
add_title(s, "Усвоение данных: Optimal Interpolation")
add_text(s, "Принцип: Учёт пространственной структуры ошибок. Информация от одиночного\nнаблюдения распространяется на окрестность через ковариационную матрицу.",
         top=Inches(1.3), font_size=20)
add_formula(s, "xₐ = x_b + K · (y_obs − H·x_b)", top=Inches(2.3), font_size=24)
add_formula(s, "K = B·Hᵀ·(H·B·Hᵀ + R)⁻¹", top=Inches(2.9), font_size=22)
add_formula(s, "B_ij = σ_b² · exp(−d_ij² / L²)", top=Inches(3.5), font_size=22)
add_bullet_list(s, [
    "R = σ_o²·I — ошибки наблюдений (σ_o — стд. откл. ошибки датчика, гиперпараметр)",
    "B — ковариация ошибок прогноза, параметризована Гауссовой ф-цией расстояния",
    "L — радиус корреляции (как далеко «доносится» инфо от датчика)",
    "σ_b — масштаб ошибки прогноза (чем больше σ_b/σ_o, тем сильнее коррекция)",
], top=Inches(4.2), font_size=18)
add_note(s, "OI «размазывает» информацию от редких датчиков на соседние области через ковариационную матрицу.", top=Inches(5.8))

# ============================================================
# SLIDE 12: OI parameter tuning
# ============================================================
s = add_slide()
add_title(s, "Подбор гиперпараметров OI")
add_text(s, "Полный перебор параметров OI (σ_b = 0.8, фиксирован). Skill Score по региону:",
         top=Inches(1.3), font_size=20)
add_table(s,
    ["Радиус (L) \\ σ_o", "0.2 (верим датчику)", "0.5 (баланс)", "0.8 (верим модели)"],
    [
        ["50 км (Локально)", "24.1%", "29.7%", "30.4%"],
        ["150 км (Средний)", "26.4%", "35.7%", "38.3%"],
        ["300 км (Широкий)", "38.7%", "**40.1%**", "40.4%"],
    ],
    top=Inches(2.3), left=Inches(1),
    col_widths=[Inches(3), Inches(2.5), Inches(2.5), Inches(3)],
    font_size=18)
add_text(s, "Выбраны: L = 300 км, σ_o = 0.5 — лучший баланс между доверием к датчикам и модели.",
         top=Inches(4.5), font_size=20, bold=True, color=GREEN)
add_note(s, "При L = 300 км информация от одной станции распространяется на ~300 км — адекватно для масштаба синоптических структур.", top=Inches(5.5))

# ============================================================
# SLIDE 13: Variable groups
# ============================================================
s = add_slide()
add_title(s, "Формирование групп переменных для усвоения")
add_text(s, "Для проверки физических гипотез мы разделили усваиваемые переменные на группы:", top=Inches(1.3), font_size=20)
add_bullet_list(s, [
    "**1. Temperature Only:** (t2m, t@850, t@500) — только температура. Достаточно ли термодинамической коррекции?",
    "**2. Surface Only:** (t2m, 10u, 10v, msl, tp) — имитация простой сети наземных метеостанций.",
    "**3. Dynamics Only:** (10u, 10v, msl, u@850, v@850, u@500, v@500, z@500) — только ветер и геопотенциал.",
    "**4. All Variables:** все переменные — эталонный режим.",
], top=Inches(2.2), font_size=20)
add_note(s, "Ключевой вопрос: достаточно ли усваивать только температуру, или модель «сдует» коррекцию неправильным ветром?", top=Inches(5.5))

# ============================================================
# SLIDE 14: Algorithm comparison by groups
# ============================================================
s = add_slide()
add_title(s, "Сравнение алгоритмов по группам переменных")
add_table(s,
    ["Группа", "Nudging (Skill)", "OI (Skill)", "Δ OI − Nudging"],
    [
        ["Temperature Only", "15.6%", "16.0%", "+0.4 п.п."],
        ["Surface Only", "16.4%", "20.0%", "**+3.6 п.п.**"],
        ["Dynamics Only", "17.8%", "**36.5%**", "**+18.7 п.п.**"],
        ["All Variables", "18.7%", "**40.1%**", "**+21.4 п.п.**"],
    ],
    top=Inches(1.5), left=Inches(1.5),
    col_widths=[Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)],
    font_size=18)
add_text(s, "Выводы:", top=Inches(4.0), font_size=20, bold=True)
add_bullet_list(s, [
    "Усваивать только температуру бесполезно — модель «сдует» коррекцию неправильным ветром.",
    "OI превосходит Nudging во всех группах, особенно для динамических переменных (+18.7 п.п.).",
    "Максимальный эффект — усвоение всех переменных с OI.",
], top=Inches(4.5), font_size=18)
add_note(s, "Коррекция полей течения (адвекции) автоматически приводит к исправлению полей температуры и влажности.", top=Inches(6.3))

# ============================================================
# SLIDE 15: Multires DA results
# ============================================================
s = add_slide()
add_title(s, "Усвоение данных на multires freeze6")
add_text(s, "Пилотный тест DA на freeze6: 8 samples, 45 узлов ROI, горизонты +6…+24ч.",
         top=Inches(1.2), font_size=18)
add_table(s,
    ["Метод", "+6ч", "+12ч", "+18ч", "+24ч", "Среднее t2m RMSE"],
    [
        ["freeze6 baseline", "0.56", "0.80", "0.78", "0.73", "0.73°C"],
        ["Nudging all", "0.53", "0.75", "0.73", "0.68", "0.67°C"],
        ["OI dyn", "0.56", "0.80", "0.78", "0.73", "0.72°C"],
        ["**OI all / surf / temp**", "**0.28**", "**0.37**", "**0.41**", "**0.43**", "**0.37°C**"],
    ],
    top=Inches(1.8), left=Inches(0.8),
    col_widths=[Inches(3.2), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(2.2)],
    font_size=16)
add_text(s, "Выводы:", top=Inches(4.65), font_size=20, bold=True)
add_bullet_list(s, [
    "Даже на short-range DA даёт сильный выигрыш: 0.73 → 0.43°C на +24ч.",
    "Лучший результат даёт OI, а не nudging.",
    "Коррекция только динамики уступает коррекции всех полей.",
], top=Inches(5.05), font_size=18)
add_note(s, "Это отдельный pilot-run на 8 samples; DA на nofreeze пока не запускался.", top=Inches(6.3))

# ============================================================
# SLIDE 16: Final result old model
# ============================================================
s = add_slide()
add_title(s, "Итоговый результат: Baseline vs OI (+24ч)")
add_table(s,
    ["Конфигурация", "RMSE t2m (°C)", "Skill (регион)"],
    [
        ["Без DA, без сшивания", "2.21", "15.5%"],
        ["+ Сшивание границ", "—", "31.5%"],
        ["+ Nudging (α=0.5, все переменные)", "—", "18.7%"],
        ["**OI + сшивание (все переменные)**", "**1.37**", "**40.1%**"],
    ],
    top=Inches(1.5), left=Inches(1.5),
    col_widths=[Inches(5), Inches(2.5), Inches(2.5)],
    font_size=18)
add_note(s, "На ранней модели (перенос 64×32 на региональную сетку) комбинация OI + сшивание даёт снижение RMSE на 38%, рост Skill на +24.6 п.п.", top=Inches(4.5))

# ============================================================
# SLIDE 17: Multires vs old
# ============================================================
s = add_slide()
add_title(s, "Мультирезолюционная модель: сравнение подходов")
add_table(s,
    ["Подход", "t2m RMSE +24ч (°C)", "Skill (регион)"],
    [
        ["Перенос 64×32 (без DA)", "2.21", "15.5%"],
        ["Перенос 64×32 + OI + сшивание", "1.37", "40.1%"],
        ["**Multires freeze6 (без DA)**", "**1.40**", "**75.8%**"],
        ["Multires nofreeze (без DA)", "1.82", "74.5%"],
        ["Multires freeze6 + DA (pilot, +24ч)", "**0.43**", "—"],
    ],
    top=Inches(1.5), left=Inches(1.5),
    col_widths=[Inches(4.5), Inches(3), Inches(2.5)],
    font_size=18)
add_text(s, "Главный результат:", top=Inches(4.5), font_size=22, bold=True, color=GREEN)
add_text(s, "Мультирезолюционный подход без усвоения уже достигает RMSE,\nсопоставимого со старой моделью с OI (1.40 vs 1.37 °C),\nпри значительно более высоком Skill (75.8% vs 40.1%).",
         top=Inches(5.0), font_size=20, color=DARK)

# ============================================================
    # SLIDE 18: Long-horizon results
# ============================================================
s = add_slide()
add_title(s, "Долгосрочный горизонт: до 9 суток")
add_text(s, "36 шагов × 6ч = 216ч (9 суток), без DA. Регион Красноярска (45 узлов), 200 сэмплов.",
         top=Inches(1.2), font_size=18)
add_table(s,
    ["Горизонт", "freeze6 t2m °C", "nofreeze t2m °C", "freeze6 Skill", "nofreeze Skill"],
    [
        ["+24ч (1 д)", "**1.37**", "1.70", "**77.3%**", "75.1%"],
        ["+48ч (2 д)", "**1.71**", "2.47", "**70.8%**", "64.8%"],
        ["+72ч (3 д)", "**2.01**", "3.32", "**61.1%**", "52.6%"],
        ["+96ч (4 д)", "**3.06**", "4.20", "**40.9%**", "38.2%"],
        ["+120ч (5 д)", "4.31", "**5.26**", "9.1%", "**16.8%**"],
        ["+144ч (6 д)", "7.51", "**5.26**", "−18.6%", "**0.3%**"],
        ["+168ч (7 д)", "10.85", "**6.83**", "−56.0%", "**−18.5%**"],
        ["+216ч (9 д)", "23.11", "**15.12**", "−144.6%", "**−94.8%**"],
    ],
    top=Inches(1.7), left=Inches(1),
    col_widths=[Inches(2), Inches(2.3), Inches(2.3), Inches(2.3), Inches(2.3)],
    font_size=16)
add_bullet_list(s, [
    "freeze6 выигрывает до 4 суток (на +24ч разница 0.33°C).",
    "После 5 суток — nofreeze деградирует плавнее (на +216ч: 15°C vs 23°C).",
    "После 5 суток обе модели становятся неустойчивыми; freeze6 деградирует быстрее.",
], top=Inches(5.8), font_size=18)
add_note(s, "Таблица — основная количественная оценка по 200 samples. Plot-bundles на 4 samples используются только как качественная иллюстрация.", top=Inches(6.55))

# ============================================================
# SLIDE 19: WRF placeholder
# ============================================================
s = add_slide()
add_title(s, "Сравнение с WRF")
add_text(s, "Физическая модель WRF (d03, ~3 км) запущена для Красноярска\nна эпизод январь 2023 («чёрное небо»).",
         top=Inches(1.3), font_size=20)
add_placeholder_box(s,
    "Карты ошибок: WRF vs GNN-модель vs ERA5\n\nДомен d03: lat [55.53, 56.39], lon [92.27, 93.62], 25 часовых шагов.",
    top=Inches(2.5))
add_text(s, "Цель сравнения:", top=Inches(5.3), font_size=20, bold=True)
add_bullet_list(s, [
    "Показать, что GNN-модель конкурирует с физической WRF при значительно меньших вычислительных затратах.",
    "Время инференса GNN: ~секунды на GPU vs часы для WRF.",
], top=Inches(5.8), font_size=18)

# ============================================================
# SLIDE 20: Live forecast with Open-Meteo comparison
# ============================================================
s = add_slide()
add_title(s, "Живой прогноз: GDAS → Open-Meteo верификация")
add_text(s, "GDAS (NOAA) → интерполяция на 133K узлов → инференс 3 суток → сравнение с Open-Meteo. Запуск 24.03.2026.",
         top=Inches(1.2), font_size=18)
add_table(s,
    ["Горизонт", "Модель T°C", "Open-Meteo T°C", "ΔT°C", "Модель MSL гПа", "Open-Meteo MSL", "ΔMSL"],
    [
        ["+6ч", "−1.67", "−1.00", "−0.67", "1008.0", "1007.5", "+0.5"],
        ["+12ч", "0.02", "0.70", "−0.68", "1009.0", "1009.0", "0.0"],
        ["+24ч", "−2.56", "−1.90", "−0.66", "1021.2", "1020.1", "+1.1"],
        ["+36ч", "1.19", "0.80", "+0.39", "1019.4", "1020.4", "−1.0"],
        ["+48ч", "3.65", "3.30", "+0.35", "1015.8", "1017.6", "−1.8"],
        ["+60ч", "5.64", "5.70", "−0.06", "1024.3", "1022.7", "+1.6"],
        ["+72ч", "2.15", "1.90", "+0.25", "1019.2", "1018.7", "+0.5"],
    ],
    top=Inches(1.7), left=Inches(0.5),
    col_widths=[Inches(1.5), Inches(1.8), Inches(1.8), Inches(1.5), Inches(2), Inches(2), Inches(1.5)],
    font_size=15)
add_text(s, "MAE температуры ≈ 0.6°C, MAE давления ≈ 1.0 гПа на горизонте до 3 суток.",
         top=Inches(5.5), font_size=20, bold=True, color=GREEN)
add_note(s, "Полный цикл «загрузка → прогноз → визуализация» занимает ~2 минуты на GPU.", top=Inches(6.2))

# ============================================================
# SLIDE 21: Conclusion
# ============================================================
s = add_slide()
add_title(s, "Заключение")
add_bullet_list(s, [
    "1. Реализована глобальная GNN-модель прогноза погоды (512×256, 19 переменных, Skill 57%).",
    "",
    "2. Предложен мультирезолюционный подход: встраивание регионального графа в глобальный — устраняет граничные артефакты и сохраняет глобальный контекст.",
    "",
    "3. Стратегия заморозки процессора снижает RMSE t2m на +24ч на 0.42°C по сравнению с обычным дообучением.",
    "",
    "4. На multires freeze6 pilot-тест DA снижает t2m RMSE на +24ч с 0.73 до 0.43°C; лучший метод — OI.",
    "",
    "5. Мультирезолюционная модель достигает Skill 75.8% (t2m RMSE = 1.40°C на +24ч) — на уровне старой модели с DA, но без усвоения.",
], top=Inches(1.3), font_size=20)
add_text(s, "Перспективы:", top=Inches(5.5), font_size=20, bold=True)
add_bullet_list(s, [
    "Усвоение данных на мультирезолюционной модели (в работе).",
    "Стабилизация авторегрессионного прогноза на горизонтах > 5 суток.",
    "Оперативный прогноз на данных GDAS.",
], top=Inches(5.9), font_size=18)

# ============================================================
# SLIDE 22: Thank you
# ============================================================
s = add_slide()
add_text(s, "Спасибо за внимание!",
         top=Inches(2.5), font_size=48, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_text(s, "Табаков Артур Станиславович",
         top=Inches(4.0), font_size=24, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_text(s, "Научный руководитель: Пененко А.В., д.ф.-м.н., ВНС ИВМиМГ СО РАН",
         top=Inches(4.7), font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
out_path = "slides.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
