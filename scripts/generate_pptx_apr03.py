"""Generate presentation PPTX — GraphCast-lite, updated 03.04.2026."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
GREEN = RGBColor(0x0B, 0x80, 0x43)
RED = RGBColor(0xCC, 0x00, 0x00)
TABLE_HEADER_BG = RGBColor(0x1A, 0x73, 0xE8)
TABLE_ALT_BG = RGBColor(0xF0, 0xF4, 0xF9)

FONT_NAME = "Times New Roman"


def _set_font(run, size=20, bold=False, color=DARK, name=FONT_NAME):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    return slide


def add_title(slide, text, top=Inches(0.3), font_size=36):
    box = slide.shapes.add_textbox(Inches(0.7), top, prs.slide_width - Inches(1.4), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, font_size, bold=True, color=BLACK)


def add_text(slide, text, top, left=Inches(0.7), width=None, font_size=20, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    if width is None:
        width = prs.slide_width - Inches(1.4)
    box = slide.shapes.add_textbox(left, top, width, Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        clean = line.replace('**', '')
        run = p.add_run()
        run.text = clean
        is_bold = bold or ('**' in line)
        _set_font(run, font_size, is_bold, color)
        p.alignment = align


def add_bullets(slide, items, top, left=Inches(0.7), width=None, font_size=20):
    if width is None:
        width = prs.slide_width - Inches(1.4)
    box = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if ':**' in item and item.startswith('**'):
            parts = item.split(':**', 1)
            label = parts[0].replace('**', '')
            rest = parts[1].strip() if len(parts) > 1 else ''
            r1 = p.add_run()
            r1.text = label + ': '
            _set_font(r1, font_size, bold=True)
            if rest:
                r2 = p.add_run()
                r2.text = rest
                _set_font(r2, font_size)
        else:
            run = p.add_run()
            run.text = item.replace('**', '')
            _set_font(run, font_size, bold=('**' in item))
        p.space_before = Pt(8)


def add_table(slide, headers, rows, top, left=Inches(0.7), col_widths=None, font_size=16):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        w = int((prs.slide_width - Inches(1.4)) / n_cols)
        col_widths = [w] * n_cols
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, sum(col_widths), Inches(0.4 * n_rows))
    tbl = shape.table
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = cw
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for par in cell.text_frame.paragraphs:
            par.font.size = Pt(font_size)
            par.font.bold = True
            par.font.color.rgb = WHITE
            par.font.name = FONT_NAME
            par.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            clean = str(val).replace('**', '')
            cell.text = clean
            for par in cell.text_frame.paragraphs:
                par.font.size = Pt(font_size)
                par.font.color.rgb = DARK
                par.font.name = FONT_NAME
                par.alignment = PP_ALIGN.CENTER
                if '**' in str(val):
                    par.font.bold = True
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG


def add_note(slide, text, top, font_size=16):
    box = slide.shapes.add_textbox(Inches(0.9), top, prs.slide_width - Inches(1.8), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, font_size, color=GRAY)
    run.font.italic = True


def add_formula(slide, text, top, font_size=20):
    box = slide.shapes.add_textbox(Inches(0.7), top, prs.slide_width - Inches(1.4), Inches(0.6))
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, font_size, color=DARK)
    p.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 1: Title
# ============================================================
s = add_slide()
add_text(s, "Повышение точности регионального\nпрогноза погоды в графовой нейросетевой\nмодели с помощью усвоения данных",
         top=Inches(1.5), font_size=40, bold=True, color=BLACK)
add_text(s, "Подготовил: Табаков Артур Станиславович\nНаучный руководитель: Пененко Алексей Владимирович, д.ф.-м.н., ВНС ИВМиМГ СО РАН",
         top=Inches(4.5), font_size=22, color=DARK)

# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
s = add_slide()
add_title(s, "Постановка проблемы")
add_bullets(s, [
    "**Накопление ошибки:** В авторегрессионном режиме нейросеть быстро теряет точность без коррекции траектории прогноза.",
    "**Разреженность наблюдений:** Реальная сеть метеостанций покрывает лишь малую часть региона; нужно восстанавливать поля в «слепых зонах».",
    "**Граничные артефакты:** Вырезание региона из глобальной сетки приводит к потере информации на границах.",
    "**Несовпадение масштабов:** Глобальная сетка нужна для динамики, а локальная детализация — для прогноза над Красноярском.",
], top=Inches(1.4), font_size=22)

# ============================================================
# SLIDE 3: Goals
# ============================================================
s = add_slide()
add_title(s, "Цель и задачи")
add_text(s, "Цель: повысить точность краткосрочного прогноза погоды над Красноярском\nв рамках графовой нейросетевой модели.",
         top=Inches(1.3), font_size=22, bold=True, color=BLACK)
add_bullets(s, [
    "1. Реализовать глобальную GNN-модель на базе архитектуры GraphCast.",
    "2. Построить мультирезолюционную схему: глобальный граф + региональная подсетка 0.25°.",
    "3. Подобрать устойчивую стратегию дообучения с заморозкой Processor.",
    "4. Исследовать усвоение данных: Nudging и Optimal Interpolation.",
    "5. Разработать каскадный downscaler (UNet) для пространственного уточнения.",
    "6. Верифицировать модель на ERA5 и в живом прогнозе по GDAS.",
], top=Inches(2.5), font_size=20)

# ============================================================
# SLIDE 4: Architecture
# ============================================================
s = add_slide()
add_title(s, "Базовая модель: архитектура GraphCast-lite")
add_bullets(s, [
    "**Encoder:** проекция признаков с регулярной grid-сетки на икосаэдральную mesh.",
    "**Processor:** message passing на mesh, обновление рёбер и узлов через MLP.",
    "**Decoder:** обратная проекция прогноза на регулярную сетку.",
], top=Inches(1.3), font_size=20)
add_text(s, "Processor = Interaction Network (Battaglia et al., 2016):",
         top=Inches(3.2), font_size=20, bold=True)
add_bullets(s, [
    "Edge MLP: фичи двух узлов + ребра → обновлённое сообщение.",
    "Node MLP: фичи узла + агрегат входящих сообщений → обновлённое состояние.",
    "12 шагов message passing, скрытая размерность 256.",
], top=Inches(3.9), font_size=18)
add_note(s, "Выбор GNN обусловлен сферической геометрией Земли — нет искажений на полюсах, характерных для прямоугольных сеток.",
         top=Inches(5.8))

# ============================================================
# SLIDE 5: Variables & Metrics
# ============================================================
s = add_slide()
add_title(s, "Набор переменных и метрики")
add_table(s,
    ["Переменная", "Описание", "Уровень"],
    [
        ["t2m", "Температура воздуха на 2 м", "поверхность"],
        ["10u, 10v", "Зональная / меридиональная компонента ветра", "10 м"],
        ["msl, sp", "Давление на уровне моря / на поверхности", "поверхность"],
        ["tp", "Суммарные осадки", "поверхность"],
        ["tcwv", "Общее влагосодержание столба атмосферы", "интегральная"],
        ["z_surf, lsm", "Геопотенциал поверхности, маска суша/море", "статическое"],
        ["t, u, v, z, q @850", "Температура, ветер, геопотенциал, влажность", "850 гПа"],
        ["t, u, v, z, q @500", "То же", "500 гПа"],
    ],
    top=Inches(1.2), left=Inches(0.7),
    col_widths=[Inches(2.5), Inches(6), Inches(3)],
    font_size=16)
add_bullets(s, [
    "**RMSE:** среднеквадратичная ошибка в физических единицах.",
    "**Skill:** 1 − RMSE_model / RMSE_persistence.",
    "**ACC:** корреляция аномалий (1 = идеал, ≤ 0 = не лучше случайного).",
], top=Inches(5.3), font_size=18)
add_note(s, "Данные: реанализ ERA5 (ECMWF), период 2010–2021, шаг 6 часов.", top=Inches(6.5))

# ============================================================
# SLIDE 6: Global Baseline
# ============================================================
s = add_slide()
add_title(s, "Глобальный baseline: прогноз на 3 суток")
add_text(s, "Сетка 512×256 (≈ 0.7°), 19 переменных, 200 тестовых сэмплов, 12 AR-шагов.",
         top=Inches(1.1), font_size=18)
add_text(s, "Глобально по всей сетке:", top=Inches(1.55), font_size=18, bold=True)
add_table(s,
    ["Горизонт", "RMSE (норм.)", "Skill vs persistence", "ACC"],
    [
        ["+24ч", "0.245", "**62.3%**", "0.969"],
        ["+48ч", "0.378", "**50.2%**", "0.928"],
        ["+72ч", "0.518", "**34.3%**", "0.869"],
        ["Среднее +6…+72ч", "0.349", "**49.7%**", "0.939"],
    ],
    top=Inches(1.95), left=Inches(1.35),
    col_widths=[Inches(2.5), Inches(2.5), Inches(2.5), Inches(2)],
    font_size=17)
add_text(s, "В регионе Красноярска на грубой глобальной сетке:", top=Inches(4.35), font_size=18, bold=True)
add_table(s,
    ["Горизонт", "t2m RMSE", "Skill", "ACC"],
    [
        ["+24ч", "0.72°C", "71.5%", "0.860"],
        ["+48ч", "1.10°C", "57.0%", "0.744"],
        ["+72ч", "1.52°C", "41.5%", "0.565"],
    ],
    top=Inches(4.75), left=Inches(2.0),
    col_widths=[Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)],
    font_size=17)
add_note(s, "На глобальной сетке регион покрыт только 3 узлами — локальные метрики не сравнимы напрямую с multires (45 узлов).", top=Inches(6.35))

# ============================================================
# SLIDE 7: Multires approach
# ============================================================
s = add_slide()
add_title(s, "Мультирезолюционный подход")
add_text(s, "Проблема: Вырезание региона из глобальной сетки → граничные артефакты.",
         top=Inches(1.3), font_size=20, bold=True, color=RED)
add_text(s, "Решение: Встраиваем региональный граф высокого разрешения в глобальный,\nформируя единую мультирезолюционную сетку.",
         top=Inches(2.0), font_size=20, bold=True, color=BLACK)
add_bullets(s, [
    "Глобальная подсетка: 512 × 256 = 131 072 узла (шаг ~0.7°)",
    "Региональная подсетка (Красноярск): 61 × 41 = 2 501 узел (шаг 0.25°)",
    "Итого: 133 279 узлов, единый граф",
], top=Inches(3.0), font_size=20)
add_text(s, "Преимущества:", top=Inches(4.5), font_size=20, bold=True)
add_bullets(s, [
    "Нет искусственных границ — глобальные узлы обеспечивают естественный контекст.",
    "Региональные узлы дают локальную детализацию в 3× выше.",
    "Encoder/Decoder работают с объединённой сеткой, Processor — на единой mesh.",
], top=Inches(5.0), font_size=18)
add_note(s, "Зона оценки (~55°–57°N, ~90°–95°E): 45 внутренних узлов вокруг Красноярска.",
         top=Inches(6.5))

# ============================================================
# SLIDE 8: Stitching / Merge
# ============================================================
s = add_slide()
add_title(s, "Сшивание глобальных и региональных данных")
add_text(s, "При построении multires-датасета для теста и инференса используется merge-режим:\nреальные региональные поля 0.25° сшиваются с глобальным полем.",
         top=Inches(1.2), font_size=20)
add_text(s, "Принцип сшивания:", top=Inches(2.2), font_size=20, bold=True)
add_formula(s, "X(v) = X_global(v),  если  v ∉ ROI", top=Inches(2.8), font_size=22)
add_formula(s, "X(v) = X_regional(v),  если  v ∈ ROI", top=Inches(3.3), font_size=22)
add_bullets(s, [
    "X_global — глобальное поле на грубой сетке 0.7°",
    "X_regional — региональное поле на сетке 0.25°",
    "ROI — целевой регион (~55°–57°N, 90°–95°E)",
    "v — узел объединённой сетки",
], top=Inches(3.8), font_size=18)
add_note(s, "Ключевая идея: это не интерполяция, а замена грубых глобальных узлов внутри ROI реальными fine-данными.",
         top=Inches(5.8))

# ============================================================
# SLIDE 9: Freeze strategy
# ============================================================
s = add_slide()
add_title(s, "Стратегия дообучения: заморозка процессора (freeze6)")
add_text(s, "Проблема: При fine-tune Processor может «забыть» глобальную физику (catastrophic forgetting).",
         top=Inches(1.3), font_size=20, color=RED)
add_text(s, "Решение:", top=Inches(2.1), font_size=20, bold=True)
add_bullets(s, [
    "Этап 1 (6 эпох): Processor ЗАМОРОЖЕН — обучаются только Encoder и Decoder.",
    "Этап 2: Processor размораживается с пониженным learning rate (×0.1).",
], top=Inches(2.6), font_size=20)
add_table(s,
    ["Метрика", "freeze6", "nofreeze"],
    [
        ["Skill (глобально)", "**66.9%**", "65.2%"],
        ["Skill (регион Красноярска)", "**75.8%**", "74.5%"],
        ["ACC (глобально)", "**0.983**", "0.981"],
    ],
    top=Inches(3.8), left=Inches(2.5),
    col_widths=[Inches(3.5), Inches(2), Inches(2)],
    font_size=18)
add_text(s, "t2m RMSE +24ч в регионе: 1.40°C (freeze6) vs 1.82°C (nofreeze).",
         top=Inches(5.5), font_size=20, bold=True, color=GREEN)
add_note(s, "Заморозка процессора критична — без неё ошибка на +24ч растёт на 30%.", top=Inches(6.2))

# ============================================================
# SLIDE 10: DA — Nudging & OI
# ============================================================
s = add_slide()
add_title(s, "Усвоение данных: Nudging и Optimal Interpolation")
add_text(s, "Задача: скорректировать прогноз по наблюдениям в 10% узлов (имитация разреженной сети станций).",
         top=Inches(1.3), font_size=20)
add_text(s, "Nudging — локальная релаксация:", top=Inches(2.1), font_size=20, bold=True)
add_formula(s, "xₐ = x_b + α · M · (y_obs − H(x_b)),  α = 0.5", top=Inches(2.6), font_size=22)
add_text(s, "Optimal Interpolation — распространение поправки через ковариационную матрицу:", top=Inches(3.3), font_size=20, bold=True)
add_formula(s, "xₐ = x_b + K · (y_obs − H·x_b)", top=Inches(3.8), font_size=22)
add_formula(s, "K = B·Hᵀ·(H·B·Hᵀ + R)⁻¹,    B_ij = σ_b² · exp(−d² / L²)", top=Inches(4.3), font_size=20)
add_bullets(s, [
    "L = 300 км — радиус корреляции (масштаб синоптических структур)",
    "σ_o = 0.5 — ошибка наблюдений, σ_b = 0.8 — ошибка прогноза",
], top=Inches(5.0), font_size=18)
add_note(s, "OI «размазывает» информацию от редких датчиков на соседние области.", top=Inches(6.2))

# ============================================================
# SLIDE 11: Variable groups comparison
# ============================================================
s = add_slide()
add_title(s, "Сравнение алгоритмов по группам переменных")
add_table(s,
    ["Группа", "Nudging (Skill)", "OI (Skill)", "Δ OI − Nudging"],
    [
        ["Temperature Only", "15.6%", "16.0%", "+0.4 п.п."],
        ["Surface Only", "16.4%", "20.0%", "**+3.6 п.п.**"],
        ["Dynamics Only", "17.8%", "**36.5%**", "**+18.7 п.п.**"],
        ["All Variables", "18.7%", "**40.1%**", "**+21.4 п.п.**"],
    ],
    top=Inches(1.5), left=Inches(1.5),
    col_widths=[Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)],
    font_size=18)
add_text(s, "Выводы:", top=Inches(4.0), font_size=20, bold=True)
add_bullets(s, [
    "Усваивать только температуру недостаточно — модель «сдувает» коррекцию ветром.",
    "OI превосходит Nudging, особенно для динамических переменных (+18.7 п.п.).",
    "Максимальный эффект — усвоение всех переменных с OI.",
], top=Inches(4.5), font_size=18)

# ============================================================
# SLIDE 12: Multires DA results
# ============================================================
s = add_slide()
add_title(s, "Усвоение данных на multires freeze6")
add_text(s, "Пилотный тест DA: 8 сэмплов, 45 узлов ROI, горизонты +6…+24ч.",
         top=Inches(1.2), font_size=18)
add_table(s,
    ["Метод", "+6ч", "+12ч", "+18ч", "+24ч", "Средн. t2m RMSE"],
    [
        ["freeze6 baseline", "0.56", "0.80", "0.78", "0.73", "0.73°C"],
        ["Nudging all", "0.53", "0.75", "0.73", "0.68", "0.67°C"],
        ["OI dyn", "0.56", "0.80", "0.78", "0.73", "0.72°C"],
        ["**OI all / surf / temp**", "**0.28**", "**0.37**", "**0.41**", "**0.43**", "**0.37°C**"],
    ],
    top=Inches(1.8), left=Inches(0.8),
    col_widths=[Inches(3.2), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(2.2)],
    font_size=16)
add_text(s, "Главный итог:", top=Inches(4.5), font_size=22, bold=True, color=GREEN)
add_text(s, "Даже short-range DA снижает t2m RMSE на +24ч с 0.73°C до 0.43°C.",
         top=Inches(5.0), font_size=20, color=DARK)

# ============================================================
# SLIDE 13: UNet concept
# ============================================================
s = add_slide()
add_title(s, "UNet downscaler: концепция каскада")
add_text(s, "U-Net — encoder-decoder архитектура со skip connections, восстанавливающая\nлокальные детали при сохранении глобального контекста.",
         top=Inches(1.2), font_size=20)
add_text(s, "Идея каскада: GNN строит физически согласованный крупномасштабный прогноз,\nзатем UNet уточняет его на fine-grid (0.25°).",
         top=Inches(2.2), font_size=20, bold=True)
add_text(s, "Формула каскада:", top=Inches(3.2), font_size=20, bold=True)
add_formula(s, "y_fine = y_coarse↑ + F_UNet(y_coarse↑, s)", top=Inches(3.8), font_size=24)
add_bullets(s, [
    "y_coarse↑ — билинейно увеличенный coarse прогноз GNN",
    "s — статические поля (рельеф, маска суша/море)",
    "F_UNet — предсказанная поправка (residual mode)",
], top=Inches(4.5), font_size=18)
add_text(s, "Параметры UNet: 7 808 211 весов, base_filters = 64, вход 40 каналов, выход 19 полей.",
         top=Inches(5.8), font_size=18, color=GRAY)

# ============================================================
# SLIDE 14: Cascade results
# ============================================================
s = add_slide()
add_title(s, "Каскад GNN → UNet на ERA5: результаты")
add_text(s, "Offline-тест: AR=4, 50 сэмплов, без постобработки.",
         top=Inches(1.1), font_size=18)
add_table(s,
    ["Переменная", "+6ч GNN", "+6ч Cascade", "Δ", "+24ч GNN", "+24ч Cascade", "Δ"],
    [
        ["t2m", "1.55°C", "1.81°C", "−17.2%", "2.08°C", "2.19°C", "−5.3%"],
        ["tp", "0.512 мм", "0.095 мм", "**+81.5%**", "0.702 мм", "0.119 мм", "**+83.1%**"],
        ["t@850", "0.79°C", "0.85°C", "−8.1%", "1.24°C", "1.18°C", "**+4.8%**"],
    ],
    top=Inches(1.5), left=Inches(0.5),
    col_widths=[Inches(1.8), Inches(1.6), Inches(1.8), Inches(1.5), Inches(1.6), Inches(1.8), Inches(1.5)],
    font_size=16)
add_table(s,
    ["Горизонт", "Skill vs persistence", "Δ Cascade vs GNN"],
    [
        ["+6ч", "32.0%", "−59.6%"],
        ["+12ч", "52.6%", "−44.2%"],
        ["+18ч", "57.0%", "−23.4%"],
        ["+24ч", "58.0%", "−20.9%"],
    ],
    top=Inches(3.6), left=Inches(2.5),
    col_widths=[Inches(2), Inches(3), Inches(3)],
    font_size=16)
add_text(s, "Каскад сильно улучшает осадки (+83%), но по t2m хуже базового GNN. Выигрыш селективный.",
         top=Inches(5.8), font_size=20, bold=True, color=DARK)

# ============================================================
# SLIDE 14b: ERA5 urban problems
# ============================================================
s = add_slide()
add_title(s, "ERA5 в городах: проблема репрезентативности")
add_bullets(s, [
    "**Grid-box smoothing:** ячейка ERA5 (~31 км) усредняет город, реку, рельеф в одно значение.",
    "**Representativeness error:** станция — точечные условия, ERA5 — среднее ячейки (Janjić et al., 2018).",
    "**Неразрешённый UHI:** теплоёмкость материалов, антропогенное тепло, уличные каньоны (Nogueira et al., 2022).",
    "**Холодные города хуже:** зимние инверсии, отопление → bias сильнее (Brozovsky et al., 2021).",
], top=Inches(1.3), font_size=20)
add_text(s, "Красноярск — особо тяжёлый случай:",
         top=Inches(3.6), font_size=22, bold=True, color=BLACK)
add_bullets(s, [
    "Енисей не замерзает → тепловой контраст с городом",
    "Частые зимние инверсии, сложный рельеф",
    "Подтверждённый UHI по Landsat (Matuzko & Yakubailik, 2018)",
], top=Inches(4.1), font_size=19)
add_note(s, "Bias ERA5 vs WMO 29570: зимой станция теплее на 5–7°C (UHI + инверсии), летом ERA5 теплее на 3–6°C днём.",
         top=Inches(5.6))

# ============================================================
# SLIDE 14c: Postprocessing pipeline
# ============================================================
s = add_slide()
add_title(s, "Постобработка: 3-ступенчатый пайплайн")
add_text(s, "1. Lapse-rate коррекция высоты", top=Inches(1.2), font_size=22, bold=True, color=BLACK)
add_formula(s, "T_corr = T_model + Γ · (z_station − z_grid),  Γ = 6.5 °C/км", top=Inches(1.65))
add_text(s, "Поправка ±0.3–1.0°C за несовпадение высот ячейки и города.", top=Inches(2.05), font_size=18)
add_text(s, "2. Learned MOS (Model Output Statistics)", top=Inches(2.6), font_size=22, bold=True, color=BLACK)
add_bullets(s, [
    "HistGradientBoostingRegressor: предсказывает bias = T_station − T_ERA5.",
    "20 признаков: метео, временны́е, лаговые, пространственные.",
    "Обучение: ERA5 + 19 станций NOAA ISD-Lite (2016–2024), r ≤ 270 км.",
    "Тест (2024): MAE = 1.32°C.",
], top=Inches(3.0), font_size=19)
add_text(s, "3. Spatial IDW — интерполяция bias на все узлы", top=Inches(4.7), font_size=22, bold=True, color=BLACK)
add_formula(s, "bias(v) = Σ wₖ·bias(sₖ) / Σ wₖ,  wₖ = 1/d(v,sₖ)²,  r_max = 300 км", top=Inches(5.15))
add_text(s, "MOS корректирует 19 точек → IDW распространяет на все 2501 узел.", top=Inches(5.55), font_size=18)

# ============================================================
# SLIDE 15: Long horizon
# ============================================================
s = add_slide()
add_title(s, "Долгосрочный горизонт: до 7 суток (без DA)")
add_text(s, "freeze6, 200 сэмплов, регион Красноярска (45 узлов).",
         top=Inches(1.2), font_size=18)
add_table(s,
    ["Горизонт", "t2m RMSE (°C)", "Skill"],
    [
        ["+24ч (1 д)", "**1.37**", "**77.3%**"],
        ["+48ч (2 д)", "**1.71**", "**70.8%**"],
        ["+72ч (3 д)", "**2.01**", "**61.1%**"],
        ["+96ч (4 д)", "3.06", "40.9%"],
        ["+120ч (5 д)", "4.31", "9.1%"],
        ["+144ч (6 д)", "7.51", "−18.6%"],
        ["+168ч (7 д)", "10.85", "−56.0%"],
    ],
    top=Inches(1.8), left=Inches(3.0),
    col_widths=[Inches(2.5), Inches(2.5), Inches(2)],
    font_size=18)
add_note(s, "Skill деградирует после 4–5 суток; после этого прогноз становится нестабильным.", top=Inches(5.8))

# ============================================================
# SLIDE 16: Live forecast day 1-2
# ============================================================
s = add_slide()
add_title(s, "Живой прогноз 03.04.2026: сутки 1–2")
add_text(s, "GDAS → GNN (freeze6) + постобработка (lapse + MOS + IDW). Данные обновлены 03.04.2026.",
         top=Inches(1.1), font_size=18)
add_table(s,
    ["Время (LT)", "GNN+postproc", "Cascade+postproc", "Open-Meteo", "Яндекс"],
    [
        ["03.04 13:00", "3.77", "3.97", "4.3", "4"],
        ["03.04 19:00", "4.30", "4.60", "3.7", "5"],
        ["04.04 01:00", "2.07", "2.29", "2.4", "2"],
        ["04.04 07:00", "0.44", "0.66", "1.1", "3"],
        ["04.04 13:00", "6.04", "6.27", "3.9", "7"],
        ["04.04 19:00", "7.51", "7.57", "4.0", "6"],
    ],
    top=Inches(1.6), left=Inches(1.0),
    col_widths=[Inches(2), Inches(2.3), Inches(2.5), Inches(2.3), Inches(2)],
    font_size=17)
add_text(s, "На первых сутках прогноз адекватен и близок к внешним сервисам.",
         top=Inches(5.5), font_size=20, bold=True, color=GREEN)

# ============================================================
# SLIDE 17: Live forecast day 3
# ============================================================
s = add_slide()
add_title(s, "Живой прогноз 03.04.2026: сутки 3")
add_table(s,
    ["Время (LT)", "GNN+postproc", "Cascade+postproc", "Open-Meteo", "Яндекс"],
    [
        ["05.04 01:00", "1.16", "1.33", "1.1", "8"],
        ["05.04 07:00", "0.18", "0.41", "−0.1", "3"],
        ["05.04 13:00", "15.83", "16.25", "10.0", "13"],
        ["05.04 19:00", "17.77", "18.12", "11.3", "10"],
        ["06.04 01:00", "9.61", "10.12", "10.3", "5"],
        ["06.04 07:00", "9.83", "10.26", "9.7", "11"],
    ],
    top=Inches(1.3), left=Inches(1.0),
    col_widths=[Inches(2), Inches(2.3), Inches(2.5), Inches(2.3), Inches(2)],
    font_size=17)
add_text(s, "Сводка:", top=Inches(4.7), font_size=20, bold=True)
add_bullets(s, [
    "GNN vs Open-Meteo: MAE 1.77°C, bias +1.4°C",
    "Cascade vs Open-Meteo: MAE 1.86°C, bias +1.68°C",
    "На 5 апреля модель даёт более резкое потепление, чем Open-Meteo.",
    "Базовый GNN + postproc надёжнее каскада по температуре в live-сценарии.",
], top=Inches(5.1), font_size=18)

# ============================================================
# SLIDE 18: Summary table
# ============================================================
s = add_slide()
add_title(s, "Сравнение подходов: сводная таблица")
add_table(s,
    ["Подход", "t2m RMSE +24ч", "Skill (регион)"],
    [
        ["Перенос 64×32 (без DA)", "2.21°C", "15.5%"],
        ["Перенос 64×32 + OI + сшивание", "1.37°C", "40.1%"],
        ["**Multires freeze6 (без DA)**", "**1.40°C**", "**75.8%**"],
        ["Multires nofreeze (без DA)", "1.82°C", "74.5%"],
        ["Multires freeze6 + DA (pilot)", "**0.43°C**", "—"],
    ],
    top=Inches(1.5), left=Inches(1.5),
    col_widths=[Inches(5), Inches(2.5), Inches(2.5)],
    font_size=18)
add_text(s, "Мультирезолюционная модель без усвоения сопоставима по RMSE\nсо старой моделью с OI (1.40 vs 1.37°C),\nно значительно лучше по Skill (75.8% vs 40.1%).",
         top=Inches(4.5), font_size=20, bold=True, color=GREEN)

# ============================================================
# SLIDE 19: Conclusion
# ============================================================
s = add_slide()
add_title(s, "Заключение")
add_bullets(s, [
    "1. Реализована глобальная GNN-модель (512×256, 19 переменных, Skill 62%).",
    "",
    "2. Предложен мультирезолюционный подход с merge-сшиванием данных; Skill в регионе — 75.8%.",
    "",
    "3. Стратегия заморозки снижает RMSE t2m на +24ч на 0.42°C по сравнению с обычным fine-tune.",
    "",
    "4. OI превосходит Nudging; pilot DA снижает RMSE с 0.73 до 0.43°C.",
    "",
    "5. UNet-каскад полезен селективно: +83% по осадкам, но не улучшает t2m в offline.",
    "",
    "6. В live-прогнозе 03.04.2026: GNN + postproc — MAE 1.77°C vs Open-Meteo.",
], top=Inches(1.3), font_size=20)
add_text(s, "Основной резерв качества — аккуратное усвоение данных,\nа не агрессивный каскадный downscaling.",
         top=Inches(5.5), font_size=20, bold=True, color=DARK)

# ============================================================
# SLIDE 20: Thank you
# ============================================================
s = add_slide()
add_text(s, "Спасибо за внимание!",
         top=Inches(2.5), font_size=48, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s, "Табаков Артур Станиславович",
         top=Inches(4.0), font_size=24, bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text(s, "Научный руководитель: Пененко А.В., д.ф.-м.н., ВНС ИВМиМГ СО РАН",
         top=Inches(4.7), font_size=20, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
out_path = "slides_apr03_urgent.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
