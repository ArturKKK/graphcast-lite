#!/usr/bin/env python3
"""
Generate an EDITABLE PPTX from slides_final.md (Marp markdown).
Uses python-pptx with proper slide layouts, styled tables, and formatted text.
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Dimensions (16:9) ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Colors ──
C_TITLE    = RGBColor(0x1B, 0x2A, 0x4A)
C_BODY     = RGBColor(0x2D, 0x2D, 0x2D)
C_ACCENT   = RGBColor(0x1A, 0x73, 0xE8)
C_CODE     = RGBColor(0xC7, 0x25, 0x4E)
C_QUOTE    = RGBColor(0x66, 0x66, 0x66)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

C_TBL_HDR     = RGBColor(0x1B, 0x2A, 0x4A)
C_TBL_HDR_FG  = C_WHITE
C_TBL_ALT     = RGBColor(0xEE, 0xF2, 0xF7)

FONT_SANS  = "Calibri"
FONT_SERIF = "Cambria"
FONT_MONO  = "Consolas"


# ─────────────────────────── Parsing ───────────────────────────

def parse_slides(md_path: str) -> list[str]:
    text = Path(md_path).read_text(encoding="utf-8")
    text = re.sub(r"^---\n.*?\n---\n", "", text, count=1, flags=re.DOTALL)
    slides = re.split(r"\n---\n", text)
    return [s.strip() for s in slides if s.strip()]


def strip_html(line: str) -> str:
    return re.sub(r"<[^>]+>", "", line).strip()


def _clean(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    return text.strip()


# ─────────────────────────── Rich text ───────────────────────────

def _add_run(para, text, size, bold=False, italic=False, color=None, font=None):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font or FONT_SANS
    if color:
        r.font.color.rgb = color
    return r


def add_rich_text(para, text, size=14, base_bold=False, color=None):
    if color is None:
        color = C_BODY
    token_re = re.compile(r'(\*\*.*?\*\*|\*[^*]+?\*|`[^`]+?`)')
    parts = token_re.split(text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            _add_run(para, part[2:-2], size, bold=True, color=color)
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            _add_run(para, part[1:-1], size, italic=True, color=color)
        elif part.startswith("`") and part.endswith("`"):
            _add_run(para, part[1:-1], size, color=C_CODE, font=FONT_MONO)
        else:
            _add_run(para, part, size, bold=base_bold, color=color)


# ─────────────────────────── Table ───────────────────────────

def parse_md_table(lines):
    header, rows = [], []
    for i, line in enumerate(lines):
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if i == 0:
            header = cells
        elif i == 1:
            continue
        else:
            rows.append(cells)
    return header, rows


def _set_cell_margin(cell, m):
    pr = cell._tc.get_or_add_tcPr()
    pr.set(qn("a:marL"), str(int(m)))
    pr.set(qn("a:marR"), str(int(m)))
    pr.set(qn("a:marT"), str(int(m // 2)))
    pr.set(qn("a:marB"), str(int(m // 2)))


def add_table(slide, header, rows, left, top, width):
    n_rows = len(rows) + 1
    n_cols = len(header)
    fsz = 11 if n_cols <= 6 else 9
    rh = Inches(0.30) if n_rows <= 8 else Inches(0.26)
    th = rh * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, th)
    table = shape.table
    tp = table._tbl.tblPr
    tp.set("bandRow", "0")
    tp.set("firstRow", "1")

    fw = int(width * 0.30)
    rw = int((width - fw) / max(n_cols - 1, 1))
    table.columns[0].width = fw
    for j in range(1, n_cols):
        table.columns[j].width = rw

    for j, h in enumerate(header):
        c = table.cell(0, j)
        c.text = ""
        p = c.text_frame.paragraphs[0]
        _add_run(p, _clean(h), fsz, bold=True, color=C_TBL_HDR_FG)
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        c.fill.solid()
        c.fill.fore_color.rgb = C_TBL_HDR
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_cell_margin(c, Inches(0.05))

    for i, row in enumerate(rows):
        for j in range(n_cols):
            val = _clean(row[j]) if j < len(row) else ""
            c = table.cell(i + 1, j)
            c.text = ""
            p = c.text_frame.paragraphs[0]
            _add_run(p, val, fsz, color=C_BODY)
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            _set_cell_margin(c, Inches(0.05))
            if i % 2 == 1:
                c.fill.solid()
                c.fill.fore_color.rgb = C_TBL_ALT
            else:
                c.fill.background()

    return shape


# ─────────────────────────── Block classifier ───────────────────────────

def classify_content(lines):
    blocks, buf = [], []
    in_table = False
    for line in lines:
        s = line.strip()
        is_tbl = s.startswith("|") and "|" in s[1:]
        if is_tbl:
            if not in_table:
                if buf:
                    blocks.append(("text", buf))
                    buf = []
                in_table = True
            buf.append(line)
        else:
            if in_table:
                blocks.append(("table", buf))
                buf = []
                in_table = False
            buf.append(line)
    if buf:
        blocks.append(("table" if in_table else "text", buf))
    return blocks


# ─────────────────────────── Slide builders ───────────────────────────

def _detect_title(lines):
    for i, line in enumerate(lines):
        s = line.strip()
        if s.startswith("# ") and not s.startswith("## "):
            return s[2:].strip(), True, i + 1
        if s.startswith("## "):
            return s[3:].strip(), False, i + 1
    return None, False, 0


def _add_accent_bar(slide, y=0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, Inches(0.06))
    s.fill.solid()
    s.fill.fore_color.rgb = C_ACCENT
    s.line.fill.background()


def build_title_slide(prs, title, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide)
    _add_accent_bar(slide, SLIDE_H - Inches(0.06))

    tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), SLIDE_W - Inches(3), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_rich_text(p, title, size=36, base_bold=True, color=C_TITLE)

    clean_body = [strip_html(l).strip() for l in body_lines if strip_html(l).strip()]
    if clean_body:
        tb2 = slide.shapes.add_textbox(Inches(2), Inches(4.2), SLIDE_W - Inches(4), Inches(2.5))
        stf = tb2.text_frame
        stf.word_wrap = True
        for j, bl in enumerate(clean_body):
            pp = stf.paragraphs[0] if j == 0 else stf.add_paragraph()
            pp.alignment = PP_ALIGN.CENTER
            pp.space_after = Pt(6)
            add_rich_text(pp, bl, size=18, color=C_BODY)


def build_content_slide(prs, title, content_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_accent_bar(slide)

    if title:
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), SLIDE_W - Inches(1.2), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        add_rich_text(p, _clean(title), size=26, base_bold=True, color=C_TITLE)

    cleaned = [strip_html(l) for l in content_lines if strip_html(l)]
    blocks = classify_content(cleaned)
    n_tables = sum(1 for t, _ in blocks if t == "table")

    LEFT = Inches(0.6)
    WIDTH = SLIDE_W - Inches(1.2)
    y = Inches(1.05)

    for btype, blines in blocks:
        if btype == "table":
            header, rows = parse_md_table(blines)
            if header and rows:
                ts = add_table(slide, header, rows, LEFT, y, WIDTH)
                y += ts.height + Inches(0.12)
        else:
            tb = slide.shapes.add_textbox(LEFT, y, WIDTH, Inches(0.1))
            tf = tb.text_frame
            tf.word_wrap = True
            first = True
            lc = 0
            for line in blines:
                s = line.strip()
                if not s:
                    continue
                pp = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                lc += 1

                if s.startswith("### "):
                    add_rich_text(pp, s[4:], size=18, base_bold=True, color=C_TITLE)
                    pp.space_before = Pt(8)
                    pp.space_after = Pt(4)
                elif s.startswith("> "):
                    add_rich_text(pp, s[2:], size=13, color=C_QUOTE)
                    pp.space_before = Pt(4)
                elif s.startswith("$$"):
                    formula = s.replace("$$", "").strip()
                    if formula:
                        pp.alignment = PP_ALIGN.CENTER
                        _add_run(pp, formula, 13, italic=True, color=C_ACCENT, font=FONT_SERIF)
                        pp.space_before = Pt(4)
                        pp.space_after = Pt(4)
                elif s.startswith("- ") or s.startswith("• "):
                    text = s[2:]
                    pp.space_before = Pt(1)
                    pp.space_after = Pt(1)
                    _add_run(pp, "●  ", 10, color=C_ACCENT)
                    add_rich_text(pp, text, size=13 if n_tables else 14, color=C_BODY)
                elif re.match(r"^\d+\.\s", s):
                    m = re.match(r"^(\d+\.)\s(.*)$", s)
                    num, text = m.group(1), m.group(2)
                    pp.space_before = Pt(1)
                    pp.space_after = Pt(1)
                    _add_run(pp, num + " ", 14, bold=True, color=C_ACCENT)
                    add_rich_text(pp, text, size=14, color=C_BODY)
                else:
                    fsz = 13 if n_tables else 14
                    add_rich_text(pp, s, size=fsz, color=C_BODY)
                    pp.space_before = Pt(2)
                    pp.space_after = Pt(2)

            y += Pt(18) * lc + Inches(0.08)


# ─────────────────────────── Main ───────────────────────────

def main():
    root = Path(__file__).resolve().parent.parent
    md_path = root / "slides_final.md"
    out_path = root / "slides_final_editable.pptx"

    slides_md = parse_slides(str(md_path))
    print(f"Parsed {len(slides_md)} slides")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, smd in enumerate(slides_md):
        lines = smd.split("\n")
        title, is_h1, cstart = _detect_title(lines)
        if is_h1:
            build_title_slide(prs, title or "", lines[cstart:])
        else:
            build_content_slide(prs, title, lines[cstart:])
        short = (title or smd[:50]).replace("\n", " ")[:60]
        print(f"  [{i+1:2d}] {short}")

    prs.save(str(out_path))
    print(f"\nSaved: {out_path}")
    print(f"Size:  {out_path.stat().st_size / 1024:.0f} KB")


if __name__ == "__main__":
    main()
