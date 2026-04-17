"""Generate final presentation PPTX (slides_final.pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
RED = RGBColor(0xCC, 0x00, 0x00)
GREEN = RGBColor(0x0B, 0x80, 0x43)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
TABLE_HEADER_BG = RGBColor(0x1A, 0x73, 0xE8)
TABLE_ALT_BG = RGBColor(0xF0, 0xF4, 0xF9)
ORANGE = RGBColor(0xE6, 0x7E, 0x00)

FONT_TITLE = 'Times New Roman'
FONT_BODY = 'Times New Roman'


def _set_font(run, size_pt, bold=False, italic=False, color=None, font=FONT_BODY):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color


def add_slide(title_text=None):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide


def add_title(slide, text, top=Inches(0.3), font_size=34, bold=True, color=BLACK):
    left = Inches(0.7)
    width = prs.slide_width - Inches(1.4)
    height = Inches(0.9)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, font_size, bold=bold, color=color, font=FONT_TITLE)
    return txBox


def add_text(slide, text, top, left=Inches(0.7), width=None, font_size=19,
             bold=False, color=DARK, alignment=PP_ALIGN.LEFT, height=Inches(5.5)):
    if width is None:
        width = prs.slide_width - Inches(1.4)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        # Detect **label:** pattern
        if '**' in line and '**:' in line:
            import re
            parts = re.split(r'\*\*(.+?)\*\*:?\s*', line)
            first = True
            for j, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                # every other element (odd index) is the bold part
                is_bold_part = (j % 2 == 1)
                run.text = part if not is_bold_part else part + ': '
                _set_font(run, font_size, bold=(bold or is_bold_part), color=color)
        else:
            clean = line.replace('**', '')
            run = p.add_run()
            run.text = clean
            _set_font(run, font_size, bold=bold, color=color)
    return txBox


def add_bullet_list(slide, items, top, left=Inches(0.7), width=None, font_size=19, indent=Inches(0.3)):
    if width is None:
        width = prs.slide_width - Inches(1.4)
    height = Inches(5.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if not item:
            p.space_before = Pt(4)
            continue
        p.space_before = Pt(6)
        # Detect **Label:** pattern
        if '**' in item:
            import re
            parts = re.split(r'\*\*(.+?)\*\*:?\s*', item)
            for j, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                is_bold_part = (j % 2 == 1)
                run.text = part if not is_bold_part else part + ': '
                _set_font(run, font_size, bold=is_bold_part, color=DARK)
        else:
            run = p.add_run()
            run.text = item
            _set_font(run, font_size, color=DARK)
    return txBox


def add_numbered_list(slide, items, top, left=Inches(0.7), width=None, font_size=19):
    if width is None:
        width = prs.slide_width - Inches(1.4)
    height = Inches(5.2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(8)
        if '**' in item:
            import re
            parts = re.split(r'\*\*(.+?)\*\*:?\s*', item)
            for j, part in enumerate(parts):
                if not part:
                    continue
                run = p.add_run()
                is_bold_part = (j % 2 == 1)
                run.text = part if not is_bold_part else part + ': '
                _set_font(run, font_size, bold=is_bold_part, color=DARK)
        else:
            run = p.add_run()
            run.text = item
            _set_font(run, font_size, color=DARK)
    return txBox


def add_table(slide, headers, rows, top, left=Inches(0.7), col_widths=None, font_size=16):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        total_w = prs.slide_width - Inches(1.4)
        col_w = int(total_w / n_cols)
        col_widths = [col_w] * n_cols

    table_width = sum(col_widths)
    table_height = Inches(0.38 * n_rows)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, table_height)
    table = shape.table

    for j, w in enumerate(col_widths):
        table.columns[j].width = w

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.name = FONT_BODY
            paragraph.font.color.rgb = WHITE
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            raw = str(val)
            clean = raw.replace('**', '')
            cell.text = clean
            is_bold = '**' in raw
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = FONT_BODY
                paragraph.font.color.rgb = DARK
                paragraph.font.bold = is_bold
                paragraph.alignment = PP_ALIGN.CENTER
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG

    return shape


def add_note(slide, text, top, font_size=15):
    left = Inches(0.9)
    width = prs.slide_width - Inches(1.8)
    height = Inches(0.9)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, font_size, italic=True, color=GRAY)
    return txBox


def add_formula(slide, text, top, left=Inches(0.7), font_size=20, center=True):
    width = prs.slide_width - Inches(1.4)
    height = Inches(0.65)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    _set_font(run, font_size, color=DARK)
    return txBox


def add_divider(slide, top):
    from pptx.util import Emu
    left = Inches(0.7)
    width = prs.slide_width - Inches(1.4)
    line = slide.shapes.add_shape(1, left, top, width, Pt(1))  # MSO_SHAPE_TYPE.LINE doesn't work like that
    # Use a thin rectangle instead
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()
    return line


# ============================================================
# SLIDE 1: Title
# ============================================================
s = add_slide()
add_text(s, "Повышение точности регионального\nпрогноза погоды в графовой нейросетевой\nмодели с помощью усвоения данных",
         top=Inches(1.4), font_size=38, bold=True, color=BLACK, alignment=PP_ALIGN.LEFT)
add_text(s, "\nПодготовил: Табаков Артур Станиславович",
         top=Inches(4.6), font_size=22, color=DARK)
add_text(s, "Научный руководитель: Пененко А.В., д.ф.-м.н., ВНС ИВМиМГ СО РАН",
         top=Inches(5.1), font_size=20, color=GRAY)

# ============================================================
# SLIDE 2: Problem Statement
# ============================================================
s = add_slide()
add_title(s, "Постановка проблемы")
add_numbered_list(s, [
    "**Накопление ошибки:** В авторегрессионном режиме нейросеть быстро теряет точность без коррекции траектории прогноза.",
    "**Разреженность наблюдений:** Реальная сеть станций покрывает малую часть региона — нужно восстанавливать поля в «слепых зонах».",
    "**Граничные артефакты:** При вырезании региона из глобальной сетки модель теряет информацию о соседях за границей ROI.",
    "**Несовпадение масштабов:** Глобальная сетка необходима для крупномасштабной динамики, а высокое разрешение — для локального прогноза.",
], top=Inches(1.35), font_size=21)

# ============================================================
# SLIDE 3: Goals
# ============================================================
s = add_slide()
add_title(s, "Цель и задачи")
add_text(s, "Цель: повысить точность краткосрочного прогноза погоды над Красноярском\nв рамках графовой нейросетевой модели.",
         top=Inches(1.2), font_size=21, bold=True, color=BLACK)
add_numbered_list(s, [
    "Реализовать глобальную GNN-модель на базе GraphCast.",
    "Построить мультирезолюционную схему: глобальный граф + региональная подсетка 0.25°.",
    "Подобрать устойчивую стратегию fine-tune с заморозкой Processor.",
    "Исследовать усвоение данных: Nudging и Optimal Interpolation.",
    "Добавить каскадный downscaler на базе UNet.",
    "Проверить модель на ERA5 и в живом прогнозе по GDAS.",
], top=Inches(2.5), font_size=20)

# ============================================================
# SLIDE 4: Architecture
# ============================================================
s = add_slide()
add_title(s, "Базовая модель: GraphCast-lite")
add_bullet_list(s, [
    "**Encoder:** проекция признаков с регулярной grid-сетки на икосаэдральную mesh.",
    "**Processor:** message passing на mesh — обновление рёбер и узлов через MLP.",
    "**Decoder:** обратная проекция прогноза на регулярную сетку.",
], top=Inches(1.3), font_size=21)

add_text(s, "Processor = Interaction Network (Battaglia et al., 2016):",
         top=Inches(3.1), font_size=20, bold=True)
add_bullet_list(s, [
    "Edge MLP: фичи двух узлов + ребро → сообщение.",
    "Node MLP: фичи узла + агрегат сообщений → обновлённое состояние.",
    "12 шагов message passing, скрытая размерность 256.",
], top=Inches(3.7), font_size=19)
add_note(s, "Выбор GNN: сферическая геометрия Земли обрабатывается естественно, без искажений полярных регионов.", top=Inches(6.4))

# ============================================================
# SLIDE 5: Variables & Metrics
# ============================================================
s = add_slide()
add_title(s, "Набор переменных и метрики")
add_text(s, "19 полей ERA5: t2m, 10u, 10v, msl, tp, sp, tcwv, z_surf, lsm; t/u/v/z/q на 850 и 500 гПа.",
         top=Inches(1.2), font_size=19)
add_table(s,
    ["Метрика", "Определение", "Интерпретация"],
    [
        ["RMSE", "Среднеквадратичная ошибка", "физические единицы (°C, м/с)"],
        ["RMSE (норм.)", "RMSE после z-score нормализации", "безразмерная, для сравнения переменных"],
        ["Skill", "1 − RMSE_model / RMSE_persistence", ">0 — лучше persistence, 100% — идеал"],
        ["ACC", "Корреляция аномалий", "1 — идеал, ≤0 — хуже случайного"],
        ["Persistence", "«Завтра = сегодня»", "наивный бейзлайн для сравнения"],
    ],
    top=Inches(1.85), left=Inches(0.7),
    col_widths=[Inches(2.5), Inches(4.5), Inches(5.0)],
    font_size=17)
add_note(s, "ERA5 (ECMWF) используется как ground truth для offline-оценки; live-режим — GDAS (NOAA).", top=Inches(5.7))

# ============================================================
# SLIDE 6: Global Baseline
# ============================================================
s = add_slide()
add_title(s, "Глобальный baseline: прогноз на 3 суток")
add_text(s, "Сетка 512×256 (~0.7°), 19 переменных, 12 AR-шагов, 200 тестовых сэмплов.",
         top=Inches(1.1), font_size=18)
add_text(s, "Глобально по всей сетке:", top=Inches(1.55), font_size=18, bold=True)
add_table(s,
    ["Горизонт", "RMSE (норм.)", "Skill vs persistence", "ACC"],
    [
        ["+24ч", "0.245", "**62.3%**", "0.969"],
        ["+48ч", "0.378", "**50.2%**", "0.928"],
        ["+72ч", "0.518", "**34.3%**", "0.869"],
        ["Среднее +6…+72ч", "0.349", "**49.7%**", "0.939"],
    ],
    top=Inches(2.0), left=Inches(1.2),
    col_widths=[Inches(2.2), Inches(2.5), Inches(2.8), Inches(2.2)],
    font_size=17)
add_text(s, "В регионе Красноярска (на грубой глобальной сетке):", top=Inches(4.35), font_size=18, bold=True)
add_table(s,
    ["Горизонт", "t2m RMSE", "Skill", "ACC"],
    [
        ["+24ч", "0.72°C", "71.5%", "0.860"],
        ["+48ч", "1.10°C", "57.0%", "0.744"],
        ["+72ч", "1.52°C", "41.5%", "0.565"],
    ],
    top=Inches(4.8), left=Inches(2.0),
    col_widths=[Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.2)],
    font_size=17)
add_note(s, "Локальные метрики в регионе нельзя напрямую сравнивать с multires: здесь всего 3 узла против 45 в multires.", top=Inches(6.4))

# ============================================================
# SLIDE 7: Multiresolution Approach
# ============================================================
s = add_slide()
add_title(s, "Мультирезолюционный подход")
add_text(s, "Проблема: вырезание региона создаёт искусственные границы → артефакты прогноза.",
         top=Inches(1.25), font_size=20, bold=True, color=RED)
add_text(s, "Решение: глобальный и региональный граф объединяются в единую плоскую сетку.",
         top=Inches(1.85), font_size=20, bold=True, color=BLACK)
add_table(s,
    ["Подсетка", "Разрешение", "Узлов"],
    [
        ["Глобальная", "~0.7° (512×256)", "131 072"],
        ["Региональная (Красноярск)", "0.25° (61×41)", "2 501"],
        ["**Итого**", "—", "**133 279**"],
    ],
    top=Inches(2.7), left=Inches(1.5),
    col_widths=[Inches(4), Inches(3), Inches(2.5)],
    font_size=18)
add_bullet_list(s, [
    "Нет искусственных границ — глобальные узлы обеспечивают контекст.",
    "Региональные узлы дают 3× более высокую детализацию в ROI.",
    "Encoder/Decoder работают с объединённой сеткой, Processor — на единой mesh.",
], top=Inches(4.65), font_size=20)
add_note(s, "Зона оценки метрик: 45 внутренних узлов вокруг Красноярска (~55°–57°N, 90°–95°E).", top=Inches(6.45))

# ============================================================
# SLIDE 8: Stitching / Merge
# ============================================================
s = add_slide()
add_title(s, "Сшивание глобальных и региональных данных")
add_text(s, "Для тестирования и инференса датасет строится в режиме merge:\n"
             "глобальные узлы внутри ROI заменяются реальными узлами высокого разрешения.",
         top=Inches(1.25), font_size=20)
add_formula(s, "X(v) = X_global(v),   если v ∉ ROI", top=Inches(2.35), font_size=21)
add_formula(s, "X(v) = X_regional(v), если v ∈ ROI", top=Inches(2.95), font_size=21)
add_bullet_list(s, [
    "X_global — глобальное поле (~0.7°, 131 072 узла)",
    "X_regional — региональное поле высокого разрешения (0.25°, 2501 узел)",
    "ROI — целевой регион (55°–57°N, 90°–95°E)",
    "v — узел объединённой сетки",
], top=Inches(3.65), font_size=20)
add_note(s, "Это НЕ интерполяция — узлы внутри ROI заменяются реальными данными высокого разрешения.", top=Inches(6.45))

# ============================================================
# SLIDE 9: Training Strategy
# ============================================================
s = add_slide()
add_title(s, "Стратегия дообучения: заморозка Processor")
add_text(s, "Проблема: при fine-tune Processor «забывает» глобальную физику (catastrophic forgetting).",
         top=Inches(1.25), font_size=20, color=RED)
add_text(s, "Двухэтапное решение — freeze → fine-tune:", top=Inches(2.0), font_size=20, bold=True)
add_bullet_list(s, [
    "Этап 1 (6 эпох): Processor ЗАМОРОЖЕН. Обучаются только Encoder и Decoder — адаптация проекции без потери глобальных знаний.",
    "Этап 2 (оставшиеся эпохи): Processor разморожен с пониженным LR (×0.1) — тонкая подстройка.",
], top=Inches(2.65), font_size=20)
add_table(s,
    ["Параметр", "Значение"],
    [
        ["freeze_epochs", "6"],
        ["LR (Encoder/Decoder)", "1e-4"],
        ["LR (Processor, этап 2)", "1e-5"],
        ["Всего эпох", "32"],
        ["AR шагов при обучении", "4"],
    ],
    top=Inches(4.4), left=Inches(3.5),
    col_widths=[Inches(3.2), Inches(2.8)],
    font_size=18)

# ============================================================
# SLIDE 10: freeze6 vs nofreeze
# ============================================================
s = add_slide()
add_title(s, "Оценка модели: freeze6 vs nofreeze")
add_text(s, "200 тестовых сэмплов, 45 узлов ROI, Красноярск.", top=Inches(1.2), font_size=18)
add_table(s,
    ["Метрика", "freeze6", "nofreeze"],
    [
        ["Skill (глобально)", "**66.9%**", "65.2%"],
        ["Skill (регион Красноярска)", "**75.8%**", "74.5%"],
        ["ACC (глобально)", "**0.983**", "0.981"],
    ],
    top=Inches(1.8), left=Inches(2.3),
    col_widths=[Inches(4.0), Inches(2.2), Inches(2.2)],
    font_size=18)
add_text(s, "Региональная RMSE по t2m (°C):", top=Inches(3.7), font_size=20, bold=True)
add_table(s,
    ["Горизонт", "freeze6", "nofreeze", "Δ"],
    [
        ["+6ч", "0.96", "0.98", "−0.02"],
        ["+12ч", "1.22", "1.33", "−0.11"],
        ["+18ч", "1.29", "1.60", "−0.31"],
        ["+24ч", "**1.40**", "1.82", "**−0.42**"],
    ],
    top=Inches(4.35), left=Inches(2.8),
    col_widths=[Inches(2.2), Inches(2.0), Inches(2.0), Inches(1.8)],
    font_size=18)
add_note(s, "Вывод: заморозка Processor критична — без неё ошибка на +24ч растёт на 30%.", top=Inches(6.45))

# ============================================================
# SLIDE 11: Long-horizon without DA
# ============================================================
s = add_slide()
add_title(s, "Долгосрочный горизонт: до 7 суток")
add_text(s, "freeze6, без DA. 200 сэмплов. Регион Красноярска.", top=Inches(1.2), font_size=18)
add_table(s,
    ["Горизонт", "freeze6 t2m °C", "Skill"],
    [
        ["+24ч (1 сут)", "1.37", "**77.3%**"],
        ["+48ч (2 сут)", "1.71", "**70.8%**"],
        ["+72ч (3 сут)", "2.01", "61.1%"],
        ["+96ч (4 сут)", "3.06", "40.9%"],
        ["+120ч (5 сут)", "4.31", "9.1%"],
        ["+144ч (6 сут)", "7.51", "−18.6%"],
        ["+168ч (7 сут)", "10.85", "−56.0%"],
    ],
    top=Inches(1.8), left=Inches(2.5),
    col_widths=[Inches(2.5), Inches(2.5), Inches(2.5)],
    font_size=18)
add_bullet_list(s, [
    "Модель устойчива до 4 суток.",
    "После 5 суток прогноз становится нестабильным без DA.",
], top=Inches(5.8), font_size=19)

# ============================================================
# SLIDE 12: Nudging
# ============================================================
s = add_slide()
add_title(s, "Усвоение данных: Nudging")
add_text(s, "Задача: скорректировать прогноз, имея наблюдения лишь в части узлов\n(имитация разреженной сети метеостанций).",
         top=Inches(1.25), font_size=20)
add_text(s, "Принцип — итеративная релаксация к наблюдению:", top=Inches(2.3), font_size=20, bold=True)
add_formula(s, "x_a = x_b + α · M · (y_obs − H(x_b))",
            top=Inches(3.05), font_size=24)
add_bullet_list(s, [
    "x_b — фоновый прогноз GNN",
    "y_obs — вектор наблюдений",
    "M — маска доступности данных (1 в точках станций, 0 вне)",
    "α — коэффициент релаксации (оптимально α = 0.5)",
], top=Inches(4.0), font_size=19)
add_note(s, "Nudging выполняется на каждом AR-шаге перед подачей данных в модель.", top=Inches(6.45))

# ============================================================
# SLIDE 13: OI
# ============================================================
s = add_slide()
add_title(s, "Усвоение данных: Optimal Interpolation")
add_text(s, "Принцип: информация от одиночного наблюдения распространяется\nна окрестность через ковариационную матрицу.",
         top=Inches(1.2), font_size=20)
add_formula(s, "x_a = x_b + K(y_obs − H·x_b)", top=Inches(2.2), font_size=23)
add_formula(s, "K = B·Hᵀ·(H·B·Hᵀ + R)⁻¹", top=Inches(2.85), font_size=21)
add_formula(s, "B_ij = σ_b² · exp(−d_ij² / L²)", top=Inches(3.45), font_size=21)
add_bullet_list(s, [
    "R = σ_o²·I — ковариация ошибок наблюдений",
    "B — ковариация ошибок прогноза (Гауссова по расстоянию)",
    "L — радиус пространственной корреляции",
    "σ_b — масштаб ошибки прогноза",
], top=Inches(4.25), font_size=19)
add_note(s, "OI «размазывает» информацию от редких датчиков на соседние области.", top=Inches(6.45))

# ============================================================
# SLIDE 14: OI hyperparameter tuning
# ============================================================
s = add_slide()
add_title(s, "Подбор гиперпараметров OI")
add_text(s, "Полный перебор. Skill Score по региону (σ_b = 0.8, фиксирован):",
         top=Inches(1.2), font_size=19)
add_table(s,
    ["Радиус L \\ σ_o", "0.2 (верим датчику)", "0.5 (баланс)", "0.8 (верим модели)"],
    [
        ["50 км (Локальный)", "24.1%", "29.7%", "30.4%"],
        ["150 км (Средний)", "26.4%", "35.7%", "38.3%"],
        ["300 км (Широкий)", "38.7%", "**40.1%**", "40.4%"],
    ],
    top=Inches(1.9), left=Inches(0.8),
    col_widths=[Inches(3.0), Inches(2.8), Inches(2.8), Inches(3.2)],
    font_size=17)
add_text(s, "Выбраны: L = 300 км, σ_o = 0.5 — лучший баланс доверия к датчикам и модели.",
         top=Inches(4.5), font_size=20, bold=True, color=GREEN)
add_note(s, "При L = 300 км информация от одной станции охватывает ~300 км — масштаб синоптических структур.", top=Inches(5.5))

# ============================================================
# SLIDE 15: Variable groups
# ============================================================
s = add_slide()
add_title(s, "Группы переменных для усвоения")
add_text(s, "Для проверки физических гипотез переменные разделены на группы:", top=Inches(1.2), font_size=20)
add_bullet_list(s, [
    "**Temperature Only:** t2m, t@850, t@500 — только температура. Достаточно ли термодинамической коррекции?",
    "**Surface Only:** t2m, 10u, 10v, msl, tp — имитация наземной метеостанции.",
    "**Dynamics Only:** 10u, 10v, msl, u/v@850, u/v/z@500 — только ветер и геопотенциал.",
    "**All Variables:** все 19 переменных — эталонный режим.",
], top=Inches(1.95), font_size=20)
add_note(s, "Ключевой вопрос: достаточно ли усваивать температуру, или модель «сдует» коррекцию неправильным ветром?", top=Inches(6.2))

# ============================================================
# SLIDE 16: Algorithm comparison
# ============================================================
s = add_slide()
add_title(s, "Сравнение алгоритмов по группам переменных")
add_table(s,
    ["Группа", "Nudging", "OI", "Δ (OI − Nudging)"],
    [
        ["Temperature Only", "15.6%", "16.0%", "+0.4 п.п."],
        ["Surface Only", "16.4%", "20.0%", "+3.6 п.п."],
        ["Dynamics Only", "17.8%", "**36.5%**", "**+18.7 п.п.**"],
        ["All Variables", "18.7%", "**40.1%**", "**+21.4 п.п.**"],
    ],
    top=Inches(1.5), left=Inches(1.3),
    col_widths=[Inches(3), Inches(2.5), Inches(2.5), Inches(2.8)],
    font_size=18)
add_text(s, "Выводы:", top=Inches(4.0), font_size=20, bold=True)
add_bullet_list(s, [
    "Усваивать только температуру бесполезно — ветер «сносит» коррекцию.",
    "OI стабильно превосходит Nudging во всех группах.",
    "Максимальный эффект — OI по всем переменным.",
], top=Inches(4.55), font_size=19)

# ============================================================
# SLIDE 17: Multires DA pilot
# ============================================================
s = add_slide()
add_title(s, "Усвоение данных на multires freeze6")
add_text(s, "Пилотный тест: 8 сэмплов, 45 узлов ROI, горизонты +6…+24ч.",
         top=Inches(1.2), font_size=18)
add_table(s,
    ["Метод", "+6ч", "+12ч", "+18ч", "+24ч", "Ср. t2m RMSE"],
    [
        ["freeze6 baseline", "0.56", "0.80", "0.78", "0.73", "0.73°C"],
        ["Nudging all", "0.53", "0.75", "0.73", "0.68", "0.67°C"],
        ["OI dyn", "0.56", "0.80", "0.78", "0.73", "0.72°C"],
        ["**OI all / surf / temp**", "**0.28**", "**0.37**", "**0.41**", "**0.43**", "**0.37°C**"],
    ],
    top=Inches(1.8), left=Inches(0.7),
    col_widths=[Inches(3.5), Inches(1.35), Inches(1.35), Inches(1.35), Inches(1.35), Inches(2.2)],
    font_size=16)
add_text(s, "Выводы:", top=Inches(4.65), font_size=20, bold=True)
add_bullet_list(s, [
    "Даже короткое DA даёт сильный выигрыш: 0.73 → 0.43°C на +24ч (−41%).",
    "OI лучше Nudging во всех конфигурациях.",
    "Коррекция только динамики уступает коррекции всех полей.",
], top=Inches(5.1), font_size=18)

# ============================================================
# SLIDE 18: UNet concept
# ============================================================
s = add_slide()
add_title(s, "UNet downscaler: концепция и каскад")
add_text(s, "U-Net — encoder-decoder архитектура со skip connections,\n"
             "эффективно восстанавливающая локальные детали при сохранении глобального контекста.",
         top=Inches(1.2), font_size=20)
add_bullet_list(s, [
    "**Encoder:** последовательное сжатие с извлечением многомасштабных признаков.",
    "**Decoder:** восстановление разрешения с объединением через skip connections.",
    "**unet_residual=True:** сеть предсказывает поправку к coarse-полю.",
], top=Inches(2.35), font_size=20)
add_text(s, "Каскад GNN → UNet:", top=Inches(4.0), font_size=20, bold=True)
add_formula(s, "y_fine = y_coarse↑ + F_UNet(y_coarse↑, s)", top=Inches(4.6), font_size=22)
add_bullet_list(s, [
    "y_coarse↑ — билинейно увеличенный прогноз GNN",
    "s — статические поля: рельеф (z_surf) и маска суша/море (lsm)",
    "Параметров UNet: 7 808 211",
], top=Inches(5.35), font_size=19)

# ============================================================
# SLIDE 19: Cascade offline results (honest)
# ============================================================
s = add_slide()
add_title(s, "Каскад GNN → UNet: честный оффлайн тест")
add_text(s, "AR=4, 50 тестовых сэмплов, unet_residual=True. ERA5 0.25°. Без postprocessing.",
         top=Inches(1.15), font_size=18)
add_table(s,
    ["Переменная", "GNN +6ч", "Casc +6ч", "Δ%", "GNN +24ч", "Casc +24ч", "Δ%"],
    [
        ["t2m (°C)", "1.55", "1.81", "−17.2%", "2.08", "2.19", "−5.3%"],
        ["tp (мм)", "0.512", "0.095", "+81.5%", "0.702", "0.119", "+83.1%"],
        ["t@850 (°C)", "0.79", "0.85", "−8.1%", "1.24", "1.18", "+4.8%"],
        ["Mean skill vs pers.", "—", "32.0%", "—", "—", "58.0%", "—"],
    ],
    top=Inches(1.8), left=Inches(0.5),
    col_widths=[Inches(2.5), Inches(1.5), Inches(1.5), Inches(1.3), Inches(1.5), Inches(1.5), Inches(1.3)],
    font_size=16)
add_text(s, "Выводы:", top=Inches(4.4), font_size=20, bold=True)
add_bullet_list(s, [
    "Осадки (tp): каскад значительно точнее GNN (+81–83%).",
    "t@850: каскад незначительно лучше при горизонтах ≥ 12ч.",
    "t2m у земли: GNN точнее Cascade (от −5% до −17%).",
    "UNet — инструмент селективного уточнения, не универсального улучшения.",
], top=Inches(4.9), font_size=18)
add_note(s, "Offline тест без postproc (CLI predict_cascade.py не поддерживает --postproc и --mos-model).", top=Inches(6.85))

# ============================================================
# SLIDE 20: Live forecast day 1
# ============================================================
s = add_slide()
add_title(s, "Живой прогноз 03.04.2026: первые сутки")
add_text(s, "GDAS 03.04.2026 → GNN + MOS постобработка → сравнение с внешними сервисами. t2m (°C), Красноярск.",
         top=Inches(1.15), font_size=18)
add_table(s,
    ["Время LT", "GNN + postproc", "Cascade + postproc", "Open-Meteo", "Яндекс"],
    [
        ["03.04 13:00", "3.77", "3.97", "4.3", "4"],
        ["03.04 19:00", "4.30", "4.60", "3.7", "5"],
        ["04.04 01:00", "2.07", "2.29", "2.4", "2"],
        ["04.04 07:00", "0.44", "0.66", "1.1", "3"],
        ["04.04 13:00", "6.04", "6.27", "3.9", "7"],
        ["04.04 19:00", "7.51", "7.57", "4.0", "6"],
    ],
    top=Inches(1.8), left=Inches(0.7),
    col_widths=[Inches(2.2), Inches(2.5), Inches(2.8), Inches(2.0), Inches(1.6)],
    font_size=16)
add_text(s, "На первых сутках прогноз адекватный и близок к внешним сервисам.",
         top=Inches(5.35), font_size=19, bold=True, color=GREEN)

# ============================================================
# SLIDE 21: Live forecast day 2-3
# ============================================================
s = add_slide()
add_title(s, "Живой прогноз 03.04.2026: вторые–третьи сутки")
add_table(s,
    ["Время LT", "GNN + postproc", "Cascade + postproc", "Open-Meteo", "Яндекс"],
    [
        ["05.04 01:00", "1.16", "1.33", "1.1", "8"],
        ["05.04 07:00", "0.18", "0.41", "−0.1", "3"],
        ["05.04 13:00", "15.83", "16.25", "10.0", "13"],
        ["05.04 19:00", "17.77", "18.12", "11.3", "10"],
        ["06.04 01:00", "9.61", "10.12", "10.3", "5"],
        ["06.04 07:00", "9.83", "10.26", "9.7", "11"],
    ],
    top=Inches(1.3), left=Inches(0.7),
    col_widths=[Inches(2.2), Inches(2.5), Inches(2.8), Inches(2.0), Inches(1.6)],
    font_size=16)
add_table(s,
    ["Сравнение", "MAE (°C)", "Bias (°C)"],
    [
        ["GNN vs Open-Meteo", "1.77", "+1.40"],
        ["Cascade vs Open-Meteo", "1.86", "+1.68"],
        ["GNN vs Яндекс", "2.67", "+0.13"],
        ["Cascade vs Яндекс", "2.65", "+0.40"],
    ],
    top=Inches(4.8), left=Inches(2.5),
    col_widths=[Inches(3.5), Inches(2.0), Inches(2.0)],
    font_size=15)
add_note(s, "5 апреля: наша модель прогнозирует более резкое потепление, чем Open-Meteo. Базовый GNN надёжнее Cascade по t2m.", top=Inches(6.85))

# ============================================================
# SLIDE 22: Conclusions
# ============================================================
s = add_slide()
add_title(s, "Основные выводы")
add_numbered_list(s, [
    "**Multires freeze6** — основная рабочая модель: Skill 75.8%, t2m RMSE +24ч = 1.40°C.",
    "**Сшивание (merge)** обеспечивает честный multires-инференс без интерполяционных артефактов.",
    "**OI** — лучший метод DA; pilot-тест: t2m RMSE +24ч 0.73 → 0.43°C.",
    "**UNet выборочно** улучшает прогноз: осадки +83%, t2m — незначительно.",
    "**Live 03.04.2026:** GNN + postproc устойчивее Cascade по температуре над городом.",
], top=Inches(1.35), font_size=21)

add_text(s, "Ключевой тезис:",
         top=Inches(5.35), font_size=20, bold=True, color=BLUE)
add_text(s, "Мультирезолюционная GNN без DA уже даёт сильный локальный прогноз;\nдальнейший резерв качества — в аккуратном усвоении данных.",
         top=Inches(5.85), font_size=19, color=DARK)

# ============================================================
# SLIDE 23: Thank you
# ============================================================
s = add_slide()
add_text(s, "Спасибо за внимание!",
         top=Inches(2.4), font_size=48, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_text(s, "Табаков Артур Станиславович",
         top=Inches(4.0), font_size=24, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_text(s, "Научный руководитель: Пененко А.В., д.ф.-м.н., ВНС ИВМиМГ СО РАН",
         top=Inches(4.65), font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
import os
out_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "slides_final.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
